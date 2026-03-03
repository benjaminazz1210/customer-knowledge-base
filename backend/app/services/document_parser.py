import io
import logging
import os
import re
import tempfile
from dataclasses import dataclass, field
from typing import List, Optional

from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..config import config

logger = logging.getLogger("nexusai.parser")


@dataclass
class StructuredSection:
    heading_path: List[str]
    heading_level: int
    content: str
    section_type: str = "text"
    page: Optional[int] = None
    slide: Optional[int] = None


@dataclass
class ExtractedImage:
    image_id: str
    image_bytes: bytes
    mime_type: str
    context: str = ""
    page: Optional[int] = None
    slide: Optional[int] = None
    source_hint: Optional[str] = None


@dataclass
class ParsedDocument:
    full_text: str
    sections: List[StructuredSection] = field(default_factory=list)
    images: List[ExtractedImage] = field(default_factory=list)
    backend_used: str = "builtin"


class DocumentParser:
    def __init__(self, backend: Optional[str] = None):
        cfg_backend = backend or os.getenv("NEXUSAI_PARSER_BACKEND") or config.DOCUMENT_PARSER_BACKEND
        self.backend = (cfg_backend or "auto").strip().lower()

    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """Backward-compatible plain text parser used by existing tests."""
        ext = os.path.splitext(filename)[1].lower()

        if ext in [".txt", ".md"]:
            return file_content.decode("utf-8")
        if ext == ".pdf":
            return DocumentParser._parse_pdf_text(file_content)
        if ext == ".docx":
            return DocumentParser._parse_docx_text(file_content)
        if ext == ".pptx":
            return DocumentParser._parse_pptx_text(file_content)
        raise ValueError(f"Unsupported file format: {ext}")

    def parse_structured(self, file_content: bytes, filename: str) -> ParsedDocument:
        backend = self.backend

        if backend in ("auto", "unstructured"):
            try:
                parsed = self._parse_with_unstructured(file_content, filename)
                if parsed.sections:
                    logger.info("📘 Structured parse backend selected: unstructured")
                    parsed.images = self._extract_images_builtin(file_content, filename)
                    return parsed
            except Exception as exc:
                if backend == "unstructured":
                    logger.warning("⚠️ Unstructured parse failed, fallback to builtin parser: %s", exc)
                else:
                    logger.info("Unstructured unavailable in auto mode, trying next backend: %s", exc)

        if backend in ("auto", "llamaparse"):
            try:
                parsed = self._parse_with_llamaparse(file_content, filename)
                if parsed.sections:
                    logger.info("📘 Structured parse backend selected: llamaparse")
                    parsed.images = self._extract_images_builtin(file_content, filename)
                    return parsed
            except Exception as exc:
                if backend == "llamaparse":
                    logger.warning("⚠️ LlamaParse parse failed, fallback to builtin parser: %s", exc)
                else:
                    logger.info("LlamaParse unavailable in auto mode, falling back to builtin: %s", exc)

        parsed = self._parse_with_builtin_structure(file_content, filename)
        parsed.backend_used = "builtin"
        logger.info("📘 Structured parse backend selected: builtin")
        return parsed

    @staticmethod
    def _parse_pdf_text(content: bytes) -> str:
        reader = PdfReader(io.BytesIO(content))
        text = ""
        for page in reader.pages:
            text += (page.extract_text() or "") + "\n"
        return text

    @staticmethod
    def _parse_docx_text(content: bytes) -> str:
        doc = Document(io.BytesIO(content))
        return "\n".join([para.text for para in doc.paragraphs])

    @staticmethod
    def _parse_pptx_text(content: bytes) -> str:
        prs = Presentation(io.BytesIO(content))
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_parts = [f"--- 第 {i+1} 页 ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            slide_parts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_parts.append(row_text)
            slides_text.append("\n".join(slide_parts))
        return "\n\n".join(slides_text)

    @staticmethod
    def _sections_from_markdown(markdown_text: str) -> List[StructuredSection]:
        sections: List[StructuredSection] = []
        heading_stack: List[str] = []
        buf: List[str] = []

        def flush_buffer():
            if not buf:
                return
            text = "\n".join(buf).strip()
            if text:
                path = heading_stack[:] if heading_stack else ["文档正文"]
                sections.append(
                    StructuredSection(
                        heading_path=path,
                        heading_level=min(max(len(path), 1), 6),
                        content=text,
                        section_type="paragraph",
                    )
                )
            buf.clear()

        for line in (markdown_text or "").splitlines():
            h = re.match(r"^(#{1,6})\s+(.+)$", line.strip())
            if h:
                flush_buffer()
                level = len(h.group(1))
                title = h.group(2).strip()
                while len(heading_stack) >= level:
                    heading_stack.pop()
                heading_stack.append(title)
                continue
            buf.append(line)

        flush_buffer()
        if not sections and markdown_text.strip():
            sections.append(
                StructuredSection(
                    heading_path=["文档正文"],
                    heading_level=1,
                    content=markdown_text.strip(),
                )
            )
        return sections

    def _parse_with_unstructured(self, file_content: bytes, filename: str) -> ParsedDocument:
        from unstructured.partition.auto import partition  # type: ignore

        elements = partition(file=io.BytesIO(file_content), file_filename=filename)
        heading_stack: List[str] = []
        sections: List[StructuredSection] = []

        for element in elements:
            text = str(element).strip()
            if not text:
                continue
            category = (getattr(element, "category", None) or element.__class__.__name__).lower()

            if category in ("title", "header", "heading", "sectionheader"):
                level = 1
                metadata = getattr(element, "metadata", None)
                category_depth = getattr(metadata, "category_depth", None)
                if category_depth is not None:
                    try:
                        level = int(category_depth) + 1
                    except Exception:
                        level = 1
                while len(heading_stack) >= level:
                    heading_stack.pop()
                heading_stack.append(text)
                continue

            path = heading_stack[:] if heading_stack else ["文档正文"]
            sections.append(
                StructuredSection(
                    heading_path=path,
                    heading_level=min(max(len(path), 1), 6),
                    content=text,
                    section_type=category,
                )
            )

        full_text = "\n\n".join(sec.content for sec in sections)
        return ParsedDocument(full_text=full_text, sections=sections, backend_used="unstructured")

    def _parse_with_llamaparse(self, file_content: bytes, filename: str) -> ParsedDocument:
        from llama_parse import LlamaParse  # type: ignore

        api_key = (os.getenv("LLAMA_CLOUD_API_KEY") or config.LLAMA_CLOUD_API_KEY or "").strip()
        if not api_key:
            raise RuntimeError("LLAMA_CLOUD_API_KEY is required for llamaparse backend")

        ext = os.path.splitext(filename)[1].lower() or ".bin"
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name

        try:
            parser = LlamaParse(api_key=api_key, result_type="markdown")
            docs = parser.load_data(tmp_path)
            markdown_parts = []
            for doc in docs:
                text = getattr(doc, "text", None) or getattr(doc, "page_content", None) or str(doc)
                if text:
                    markdown_parts.append(text)
            markdown_text = "\n\n".join(markdown_parts).strip()
            sections = self._sections_from_markdown(markdown_text)
            full_text = "\n\n".join(sec.content for sec in sections)
            return ParsedDocument(full_text=full_text, sections=sections, backend_used="llamaparse")
        finally:
            try:
                os.remove(tmp_path)
            except Exception:
                pass

    def _parse_with_builtin_structure(self, file_content: bytes, filename: str) -> ParsedDocument:
        ext = os.path.splitext(filename)[1].lower()

        if ext in [".txt", ".md"]:
            text = file_content.decode("utf-8")
            if ext == ".md":
                sections = self._sections_from_markdown(text)
            else:
                sections = [
                    StructuredSection(
                        heading_path=["文档正文"],
                        heading_level=1,
                        content=text.strip(),
                        section_type="paragraph",
                    )
                ] if text.strip() else []
            return ParsedDocument(
                full_text=text,
                sections=sections,
                images=[],
                backend_used="builtin",
            )

        if ext == ".pdf":
            return self._parse_pdf_structured(file_content)
        if ext == ".docx":
            return self._parse_docx_structured(file_content)
        if ext == ".pptx":
            return self._parse_pptx_structured(file_content)
        raise ValueError(f"Unsupported file format: {ext}")

    def _parse_pdf_structured(self, content: bytes) -> ParsedDocument:
        reader = PdfReader(io.BytesIO(content))
        sections: List[StructuredSection] = []

        for idx, page in enumerate(reader.pages):
            text = (page.extract_text() or "").strip()
            if not text:
                continue
            sections.append(
                StructuredSection(
                    heading_path=[f"第{idx + 1}页"],
                    heading_level=1,
                    content=text,
                    section_type="page_text",
                    page=idx + 1,
                )
            )

        full_text = "\n\n".join(sec.content for sec in sections)
        images = self._extract_pdf_images(content)
        return ParsedDocument(full_text=full_text, sections=sections, images=images, backend_used="builtin")

    def _parse_docx_structured(self, content: bytes) -> ParsedDocument:
        doc = Document(io.BytesIO(content))
        sections: List[StructuredSection] = []
        heading_stack: List[str] = []

        for para in doc.paragraphs:
            text = (para.text or "").strip()
            if not text:
                continue

            style_name = (para.style.name if para.style else "").lower()
            hm = re.search(r"heading\s*(\d+)", style_name)
            if hm:
                level = max(1, min(int(hm.group(1)), 6))
                while len(heading_stack) >= level:
                    heading_stack.pop()
                heading_stack.append(text)
                continue

            path = heading_stack[:] if heading_stack else ["文档正文"]
            sections.append(
                StructuredSection(
                    heading_path=path,
                    heading_level=min(max(len(path), 1), 6),
                    content=text,
                    section_type="paragraph",
                )
            )

        full_text = "\n\n".join(sec.content for sec in sections)
        images = self._extract_docx_images(content)
        return ParsedDocument(full_text=full_text, sections=sections, images=images, backend_used="builtin")

    def _parse_pptx_structured(self, content: bytes) -> ParsedDocument:
        prs = Presentation(io.BytesIO(content))
        sections: List[StructuredSection] = []
        images: List[ExtractedImage] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_title = ""
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()
            if not slide_title:
                slide_title = f"第{slide_idx}页"

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = (para.text or "").strip()
                        if not text:
                            continue
                        level = int(getattr(para, "level", 0))
                        sections.append(
                            StructuredSection(
                                heading_path=[f"第{slide_idx}页", slide_title],
                                heading_level=min(2 + level, 6),
                                content=text,
                                section_type="slide_text",
                                slide=slide_idx,
                            )
                        )

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
                    try:
                        img_bytes = shape.image.blob
                        img_ext = (shape.image.ext or "png").lower()
                        mime = self._mime_from_ext(img_ext)
                        images.append(
                            ExtractedImage(
                                image_id=f"slide_{slide_idx}_image_{len(images)+1}",
                                image_bytes=img_bytes,
                                mime_type=mime,
                                context=slide_title,
                                slide=slide_idx,
                                source_hint=f"slide_{slide_idx}",
                            )
                        )
                    except Exception:
                        continue

        full_text = "\n\n".join(sec.content for sec in sections)
        return ParsedDocument(full_text=full_text, sections=sections, images=images, backend_used="builtin")

    def _extract_images_builtin(self, file_content: bytes, filename: str) -> List[ExtractedImage]:
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".pdf":
            return self._extract_pdf_images(file_content)
        if ext == ".docx":
            return self._extract_docx_images(file_content)
        if ext == ".pptx":
            return self._extract_pptx_images(file_content)
        return []

    def _extract_pdf_images(self, content: bytes) -> List[ExtractedImage]:
        reader = PdfReader(io.BytesIO(content))
        images: List[ExtractedImage] = []
        for page_idx, page in enumerate(reader.pages, start=1):
            page_images = getattr(page, "images", None) or []
            for i, image_file in enumerate(page_images, start=1):
                data = getattr(image_file, "data", None)
                name = getattr(image_file, "name", None) or f"image_{i}.png"
                if not data:
                    continue
                ext = (os.path.splitext(name)[1].replace(".", "") or "png").lower()
                images.append(
                    ExtractedImage(
                        image_id=f"pdf_{page_idx}_{i}",
                        image_bytes=data,
                        mime_type=self._mime_from_ext(ext),
                        context=f"PDF第{page_idx}页图片",
                        page=page_idx,
                        source_hint=name,
                    )
                )
        return images

    def _extract_docx_images(self, content: bytes) -> List[ExtractedImage]:
        doc = Document(io.BytesIO(content))
        images: List[ExtractedImage] = []
        seen_part = set()
        for rel in doc.part._rels.values():
            if "image" not in rel.reltype:
                continue
            image_part = rel.target_part
            key = str(getattr(image_part, "partname", ""))
            if key in seen_part:
                continue
            seen_part.add(key)
            blob = getattr(image_part, "blob", None)
            if not blob:
                continue
            ext = os.path.splitext(key)[1].replace(".", "").lower() or "png"
            images.append(
                ExtractedImage(
                    image_id=f"docx_{len(images)+1}",
                    image_bytes=blob,
                    mime_type=self._mime_from_ext(ext),
                    context="Word文档插图",
                    source_hint=key,
                )
            )
        return images

    def _extract_pptx_images(self, content: bytes) -> List[ExtractedImage]:
        prs = Presentation(io.BytesIO(content))
        images: List[ExtractedImage] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else f"第{slide_idx}页"
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
                    try:
                        img_bytes = shape.image.blob
                        img_ext = (shape.image.ext or "png").lower()
                        images.append(
                            ExtractedImage(
                                image_id=f"pptx_{slide_idx}_{len(images)+1}",
                                image_bytes=img_bytes,
                                mime_type=self._mime_from_ext(img_ext),
                                context=slide_title,
                                slide=slide_idx,
                                source_hint=f"slide_{slide_idx}",
                            )
                        )
                    except Exception:
                        continue
        return images

    @staticmethod
    def _mime_from_ext(ext: str) -> str:
        mapping = {
            "png": "image/png",
            "jpg": "image/jpeg",
            "jpeg": "image/jpeg",
            "bmp": "image/bmp",
            "gif": "image/gif",
            "webp": "image/webp",
            "tif": "image/tiff",
            "tiff": "image/tiff",
        }
        return mapping.get(ext.lower(), "application/octet-stream")
