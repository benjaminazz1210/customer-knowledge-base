import copy
import json
import logging
import re
import time
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.parse import urlparse

from docx import Document
from openai import OpenAI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ..config import config
from .embedding_service import EmbeddingService
from .vector_store import VectorStore

try:
    from langgraph.graph import END, START, StateGraph

    HAS_LANGGRAPH = True
except Exception:
    HAS_LANGGRAPH = False
    START = "__start__"
    END = "__end__"
    StateGraph = None  # type: ignore

logger = logging.getLogger("nexusai.workflow")


class WorkflowService:
    def __init__(self):
        backend_root = Path(__file__).resolve().parents[2]
        self.output_root = (backend_root / config.WORKFLOW_OUTPUT_DIR).resolve()
        self.jobs_dir = self.output_root / "jobs"
        self.files_dir = self.output_root / "files"
        self.template_dir = (backend_root / config.WORKFLOW_TEMPLATE_DIR).resolve()

        self.jobs_dir.mkdir(parents=True, exist_ok=True)
        self.files_dir.mkdir(parents=True, exist_ok=True)
        self.template_dir.mkdir(parents=True, exist_ok=True)

        self.embedding_service: Optional[EmbeddingService] = None
        self.vector_store: Optional[VectorStore] = None
        self.llm_client = self._build_llm_client()

        self.generate_graph = self._build_generate_graph() if HAS_LANGGRAPH else None
        self.revise_graph = self._build_revise_graph() if HAS_LANGGRAPH else None
        logger.info("🧩 WorkflowService initialized (langgraph=%s)", HAS_LANGGRAPH)

    @staticmethod
    def _normalize_base_url(base_url: str) -> str:
        normalized = (base_url or "").strip().rstrip("/")
        if not normalized:
            return normalized
        parsed = urlparse(normalized)
        if parsed.path in ("", "/"):
            return f"{normalized}/v1"
        return normalized

    def _build_llm_client(self):
        provider = config.LLM_PROVIDER.strip().lower()
        try:
            if provider == "ollama":
                return OpenAI(api_key="ollama", base_url=self._normalize_base_url(config.OLLAMA_BASE_URL))
            if provider == "deepseek" and config.DEEPSEEK_API_KEY:
                return OpenAI(
                    api_key=config.DEEPSEEK_API_KEY,
                    base_url=self._normalize_base_url(config.DEEPSEEK_BASE_URL),
                )
            if provider in ("openai", "heiyucode") and config.OPENAI_API_KEY:
                return OpenAI(
                    api_key=config.OPENAI_API_KEY,
                    base_url=self._normalize_base_url(config.OPENAI_BASE_URL),
                )
        except Exception as exc:
            logger.warning("⚠️ Failed to initialize LLM client for workflow: %s", exc)
        logger.warning("⚠️ WorkflowService using rule-based drafting fallback")
        return None

    def _ensure_retrieval_clients(self):
        if self.embedding_service is None:
            self.embedding_service = EmbeddingService()
        if self.vector_store is None:
            self.vector_store = VectorStore()

    @staticmethod
    def _slug(text: str, fallback: str = "document") -> str:
        cleaned = re.sub(r"\s+", "-", (text or "").strip())
        cleaned = re.sub(r"[^a-zA-Z0-9\-\u4e00-\u9fff]", "", cleaned)
        cleaned = cleaned.strip("-")
        return cleaned[:36] if cleaned else fallback

    @staticmethod
    def _extract_num(pattern: str, text: str) -> Optional[int]:
        match = re.search(pattern, text or "", re.IGNORECASE)
        if not match:
            return None
        try:
            return int(match.group(1))
        except Exception:
            return None

    @staticmethod
    def _contains_any(text: str, keywords: List[str]) -> bool:
        lowered = (text or "").lower()
        return any(k.lower() in lowered for k in keywords)

    def _infer_file_type(self, prompt: str, explicit_file_type: Optional[str]) -> str:
        if explicit_file_type in ("docx", "pptx"):
            return explicit_file_type
        lowered = (prompt or "").lower()
        if "ppt" in lowered or "幻灯" in prompt or "演示" in prompt:
            return "pptx"
        return "docx"

    def _build_generate_graph(self):
        graph = StateGraph(dict)
        graph.add_node("requirements_analyst", self._node_requirements_analyst)
        graph.add_node("retriever", self._node_retriever)
        graph.add_node("writer", self._node_writer)
        graph.add_node("ppt_layout", self._node_ppt_layout)
        graph.add_node("renderer", self._node_renderer)

        graph.add_edge(START, "requirements_analyst")
        graph.add_edge("requirements_analyst", "retriever")
        graph.add_edge("retriever", "writer")
        graph.add_conditional_edges(
            "writer",
            self._route_after_writer,
            {
                "ppt_layout": "ppt_layout",
                "renderer": "renderer",
            },
        )
        graph.add_edge("ppt_layout", "renderer")
        graph.add_edge("renderer", END)
        return graph.compile()

    def _build_revise_graph(self):
        graph = StateGraph(dict)
        graph.add_node("load_previous", self._node_load_previous)
        graph.add_node("hitl", self._node_human_feedback_router)
        graph.add_node("retriever", self._node_retriever)
        graph.add_node("rewrite", self._node_rewrite)
        graph.add_node("ppt_layout", self._node_ppt_layout)
        graph.add_node("renderer", self._node_renderer)

        graph.add_edge(START, "load_previous")
        graph.add_edge("load_previous", "hitl")
        graph.add_edge("hitl", "retriever")
        graph.add_edge("retriever", "rewrite")
        graph.add_conditional_edges(
            "rewrite",
            self._route_after_writer,
            {
                "ppt_layout": "ppt_layout",
                "renderer": "renderer",
            },
        )
        graph.add_edge("ppt_layout", "renderer")
        graph.add_edge("renderer", END)
        return graph.compile()

    def _route_after_writer(self, state: Dict[str, Any]) -> str:
        requirements = state.get("requirements", {}) or {}
        file_type = requirements.get("file_type", "docx")
        return "ppt_layout" if file_type == "pptx" else "renderer"

    def _node_requirements_analyst(self, state: Dict[str, Any]) -> Dict[str, Any]:
        prompt = state.get("prompt", "")
        explicit_file_type = state.get("file_type")
        file_type = self._infer_file_type(prompt, explicit_file_type)

        target_words = state.get("target_words") or self._extract_num(r"(\d{2,5})\s*字", prompt)
        target_slides = state.get("target_slides") or self._extract_num(r"(\d{1,2})\s*页", prompt)

        if file_type == "docx" and not target_words:
            target_words = 1000
        if file_type == "pptx" and not target_slides:
            target_slides = 7

        requires_kpi = self._contains_any(
            prompt,
            ["kpi", "可量化", "考核", "指标清单", "绩效"],
        )

        title_seed = re.sub(r"[\n\r]+", " ", prompt).strip()[:50]
        default_title = f"智能生成方案-{time.strftime('%Y%m%d')}"
        title = title_seed if title_seed else default_title

        state["requirements"] = {
            "file_type": file_type,
            "target_words": int(target_words) if target_words else None,
            "target_slides": int(target_slides) if target_slides else None,
            "requires_kpi": requires_kpi,
            "title": title,
            "use_rag": bool(state.get("use_rag", True)),
            "template_name": state.get("template_name"),
        }
        return state

    def _node_load_previous(self, state: Dict[str, Any]) -> Dict[str, Any]:
        job_id = state.get("job_id")
        if not job_id:
            raise ValueError("revise requires job_id")
        prev = self.get_job(job_id)
        if not prev:
            raise ValueError(f"job not found: {job_id}")
        state["previous_job"] = prev
        state["prompt"] = prev.get("prompt", "")
        state["requirements"] = prev.get("requirements", {})
        state["plan"] = prev.get("plan", {})
        state["draft"] = prev.get("draft", {})
        state["layout"] = prev.get("layout", {})
        return state

    @staticmethod
    def _parse_feedback_scope(feedback: str, file_type: str) -> Dict[str, Any]:
        feedback = feedback or ""
        cn_num_map = {"零": 0, "一": 1, "二": 2, "两": 2, "三": 3, "四": 4, "五": 5, "六": 6, "七": 7, "八": 8, "九": 9}

        def cn_to_int(raw: str) -> Optional[int]:
            if not raw:
                return None
            if raw.isdigit():
                return int(raw)
            if raw == "十":
                return 10
            if "十" in raw:
                parts = raw.split("十")
                left = cn_num_map.get(parts[0], 1) if parts[0] else 1
                right = cn_num_map.get(parts[1], 0) if len(parts) > 1 and parts[1] else 0
                return left * 10 + right
            total = 0
            for ch in raw:
                if ch in cn_num_map:
                    total = total * 10 + cn_num_map[ch]
                else:
                    return None
            return total if total > 0 else None

        if file_type == "pptx":
            m = re.search(r"第\s*(\d+)\s*页", feedback)
            if m:
                return {"type": "slide", "index": max(int(m.group(1)) - 1, 0)}
            m_cn = re.search(r"第\s*([一二三四五六七八九十两]+)\s*页", feedback)
            if m_cn:
                page = cn_to_int(m_cn.group(1))
                if page:
                    return {"type": "slide", "index": max(page - 1, 0)}
        if file_type == "docx":
            m = re.search(r"第\s*(\d+)\s*(节|部分|章|点|条)", feedback)
            if m:
                return {"type": "section", "index": max(int(m.group(1)) - 1, 0)}
            m_cn = re.search(r"第\s*([一二三四五六七八九十两]+)\s*(节|部分|章|点|条)", feedback)
            if m_cn:
                sec = cn_to_int(m_cn.group(1))
                if sec:
                    return {"type": "section", "index": max(sec - 1, 0)}
        return {"type": "global"}

    def _node_human_feedback_router(self, state: Dict[str, Any]) -> Dict[str, Any]:
        feedback = state.get("feedback", "")
        requirements = state.get("requirements", {}) or {}
        file_type = requirements.get("file_type", "docx")
        scope = self._parse_feedback_scope(feedback, file_type)

        if self._contains_any(feedback, ["kpi", "可量化", "指标清单", "考核表", "考核指标", "绩效指标"]):
            scope["inject_kpi"] = True
        state["feedback_scope"] = scope
        return state

    def _node_retriever(self, state: Dict[str, Any]) -> Dict[str, Any]:
        requirements = state.get("requirements", {}) or {}
        if not requirements.get("use_rag", True):
            state["retrieved_chunks"] = []
            return state

        query_text = state.get("feedback") or state.get("prompt") or ""
        if not query_text.strip():
            state["retrieved_chunks"] = []
            return state

        try:
            self._ensure_retrieval_clients()
            query_vector = self.embedding_service.get_embeddings([query_text])[0]  # type: ignore
            hits = self.vector_store.hybrid_search(  # type: ignore
                query_text=query_text,
                query_vector=query_vector,
                limit=max(1, int(config.WORKFLOW_MAX_CONTEXT_CHUNKS)),
                alpha=float(config.WORKFLOW_HYBRID_ALPHA),
            )
            state["retrieved_chunks"] = hits
        except Exception as exc:
            logger.warning("⚠️ Workflow retriever fallback to empty context: %s", exc)
            state["retrieved_chunks"] = []
        return state

    @staticmethod
    def _summarize_sources(retrieved_chunks: List[Dict[str, Any]], max_chars: int = 220) -> List[str]:
        snippets: List[str] = []
        for hit in retrieved_chunks or []:
            payload = hit.get("payload", {}) or {}
            source_file = payload.get("source_file", "unknown")
            raw = str(payload.get("chunk_text", "")).strip().replace("\n", " ")
            if not raw:
                continue
            snippets.append(f"[{source_file}] {raw[:max_chars]}")
        return snippets[:6]

    @staticmethod
    def _build_kpi_checklist(topic: str) -> List[Dict[str, str]]:
        topic = topic or "项目"
        return [
            {
                "name": "按期交付率",
                "formula": "按期完成任务数 / 计划任务总数",
                "target": ">= 95%",
                "cycle": "月度",
                "owner": "项目经理",
            },
            {
                "name": "问题闭环时效",
                "formula": "问题关闭平均时长（小时）",
                "target": "<= 72 小时",
                "cycle": "周度",
                "owner": "运营负责人",
            },
            {
                "name": f"{topic}覆盖率",
                "formula": "已覆盖对象数 / 应覆盖对象总数",
                "target": ">= 90%",
                "cycle": "月度",
                "owner": "业务主管",
            },
            {
                "name": "用户满意度",
                "formula": "满意问卷数 / 有效问卷数",
                "target": ">= 90%",
                "cycle": "季度",
                "owner": "服务团队",
            },
            {
                "name": "数据完整率",
                "formula": "完整字段记录数 / 总记录数",
                "target": ">= 98%",
                "cycle": "周度",
                "owner": "数据管理员",
            },
        ]

    def _llm_text(self, system_prompt: str, user_prompt: str, temperature: float = 0.4) -> Optional[str]:
        if not self.llm_client:
            return None
        try:
            resp = self.llm_client.chat.completions.create(
                model=config.LLM_MODEL,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ],
                temperature=temperature,
            )
            if not resp or not getattr(resp, "choices", None):
                return None
            return resp.choices[0].message.content
        except Exception as exc:
            logger.warning("⚠️ Workflow LLM generation failed, fallback to rule-based: %s", exc)
            return None

    @staticmethod
    def _split_sentences(text: str) -> List[str]:
        text = (text or "").strip()
        if not text:
            return []
        parts = re.split(r"[。！？\n]+", text)
        return [p.strip() for p in parts if p.strip()]

    def _build_rule_based_docx_draft(self, state: Dict[str, Any]) -> Dict[str, Any]:
        req = state.get("requirements", {}) or {}
        prompt = state.get("prompt", "")
        snippets = self._summarize_sources(state.get("retrieved_chunks", []))
        target_words = req.get("target_words") or 1000
        title = req.get("title") or "智能文档"

        sections: List[Dict[str, Any]] = []
        base_headings = ["一、需求背景", "二、现状分析", "三、实施方案", "四、推进计划", "五、风险与保障"]
        for idx, heading in enumerate(base_headings):
            source_line = snippets[idx % len(snippets)] if snippets else "暂无可用知识库证据，以下为通用方案框架。"
            paragraphs = [
                f"围绕“{prompt}”，本节从业务目标与约束条件进行梳理，明确实施边界与验收标准。",
                f"证据摘要：{source_line}",
                "建议采用“目标分解-任务编排-过程监控-结果复盘”的闭环机制，确保方案可执行与可评估。",
            ]
            sections.append({"heading": heading, "paragraphs": paragraphs})

        if req.get("requires_kpi"):
            kpis = self._build_kpi_checklist(prompt[:20] or "项目")
            sections.append(
                {
                    "heading": "六、可量化KPI指标清单（便于纳入区级考核）",
                    "paragraphs": [
                        "以下指标可直接纳入季度或年度考核：",
                        *[
                            f"{i+1}. {item['name']} | 口径: {item['formula']} | 目标: {item['target']} | 周期: {item['cycle']} | 责任: {item['owner']}"
                            for i, item in enumerate(kpis)
                        ],
                    ],
                    "kpis": kpis,
                }
            )

        approx_words = sum(len("".join(sec["paragraphs"])) for sec in sections)
        filler_needed = max(0, target_words - approx_words)
        if filler_needed > 120:
            filler_text = "执行要点：建立跨部门协同机制，设置里程碑评审，按周跟踪关键任务，按月复盘偏差并纠偏。"
            sections[-1]["paragraphs"].append(filler_text * max(1, filler_needed // max(len(filler_text), 1)))

        return {
            "title": title,
            "sections": sections,
            "word_target": target_words,
            "analysis": {
                "top_down": [
                    "目标分解",
                    "路径设计",
                    "组织保障",
                    "量化评估",
                ],
                "bottom_up": snippets,
            },
        }

    def _build_rule_based_ppt_draft(self, state: Dict[str, Any]) -> Dict[str, Any]:
        req = state.get("requirements", {}) or {}
        prompt = state.get("prompt", "")
        snippets = self._summarize_sources(state.get("retrieved_chunks", []))
        target_slides = req.get("target_slides") or 7
        title = req.get("title") or "智能方案汇报"

        slide_titles = [
            "封面",
            "项目背景与目标",
            "现状问题诊断",
            "总体方案设计",
            "实施路径与里程碑",
            "风险控制与保障",
            "可量化KPI与验收",
        ]
        slides: List[Dict[str, Any]] = []
        for i in range(target_slides):
            st = slide_titles[i] if i < len(slide_titles) else f"专题页 {i+1}"
            source_line = snippets[i % len(snippets)] if snippets else "暂无知识库证据"
            bullets = [
                f"围绕“{prompt}”定义本页核心结论与行动建议",
                f"证据参考：{source_line}",
                "明确责任人、时间节点、交付标准，形成可追踪闭环",
            ]
            slides.append(
                {
                    "page": i + 1,
                    "title": st,
                    "bullets": bullets,
                    "speaker_notes": "讲解时突出可执行性与量化指标。",
                }
            )

        if req.get("requires_kpi"):
            kpis = self._build_kpi_checklist(prompt[:20] or "项目")
            slides.append(
                {
                    "page": len(slides) + 1,
                    "title": "可量化KPI指标清单（便于纳入区级考核）",
                    "bullets": [f"{k['name']}：目标 {k['target']}，周期 {k['cycle']}" for k in kpis],
                    "speaker_notes": "此页用于考核口径对齐。",
                    "kpis": kpis,
                }
            )

        return {
            "title": title,
            "slides": slides,
            "analysis": {
                "top_down": ["目标", "机制", "资源", "评估"],
                "bottom_up": snippets,
            },
        }

    def _node_writer(self, state: Dict[str, Any]) -> Dict[str, Any]:
        req = state.get("requirements", {}) or {}
        file_type = req.get("file_type", "docx")
        context_lines = self._summarize_sources(state.get("retrieved_chunks", []))
        context_text = "\n".join(context_lines) if context_lines else "暂无检索上下文。"

        llm_draft = None
        system_prompt = (
            "你是企业知识库写作 Agent。输出要结构化、可执行、可考核，禁止编造无法验证的数据。"
        )
        if file_type == "docx":
            user_prompt = (
                f"需求：{state.get('prompt', '')}\n"
                f"目标字数：{req.get('target_words', 1000)}\n"
                f"上下文：\n{context_text}\n\n"
                "请输出纯文本，包含明确的小节标题和条理化内容。"
            )
            llm_draft = self._llm_text(system_prompt, user_prompt)
            state["draft"] = self._build_rule_based_docx_draft(state)
            if llm_draft:
                state["draft"]["llm_enhancement"] = llm_draft[:4000]
        else:
            user_prompt = (
                f"需求：{state.get('prompt', '')}\n"
                f"目标页数：{req.get('target_slides', 7)}\n"
                f"上下文：\n{context_text}\n\n"
                "请输出结构化PPT内容建议。"
            )
            llm_draft = self._llm_text(system_prompt, user_prompt)
            state["draft"] = self._build_rule_based_ppt_draft(state)
            if llm_draft:
                state["draft"]["llm_enhancement"] = llm_draft[:3000]
        return state

    def _node_rewrite(self, state: Dict[str, Any]) -> Dict[str, Any]:
        previous = state.get("previous_job", {}) or {}
        prev_draft = copy.deepcopy(previous.get("draft", {}))
        feedback = state.get("feedback", "")
        scope = state.get("feedback_scope", {}) or {"type": "global"}
        req = state.get("requirements", {}) or {}
        file_type = req.get("file_type", "docx")

        if file_type == "docx":
            sections = prev_draft.get("sections", [])
            if not isinstance(sections, list):
                sections = []
            if scope.get("type") == "section" and sections:
                idx = min(scope.get("index", 0), len(sections) - 1)
                sections[idx].setdefault("paragraphs", []).append(f"根据反馈优化：{feedback}")
            else:
                sections.append(
                    {
                        "heading": "反馈优化补充",
                        "paragraphs": [
                            f"收到用户反馈：{feedback}",
                            "已基于上下文分析进行定向补充，并确保与原文结构一致。",
                        ],
                    }
                )
            if scope.get("inject_kpi"):
                kpis = self._build_kpi_checklist(state.get("prompt", "")[:20] or "项目")
                sections.append(
                    {
                        "heading": "可量化KPI指标清单（便于纳入区级考核）",
                        "paragraphs": [
                            f"{i+1}. {k['name']} | 目标 {k['target']} | 周期 {k['cycle']}"
                            for i, k in enumerate(kpis)
                        ],
                        "kpis": kpis,
                    }
                )
            prev_draft["sections"] = sections
            prev_draft["revision_note"] = feedback
            state["draft"] = prev_draft
            return state

        slides = prev_draft.get("slides", [])
        if not isinstance(slides, list):
            slides = []
        if scope.get("type") == "slide" and slides:
            idx = min(scope.get("index", 0), len(slides) - 1)
            slides[idx].setdefault("bullets", []).append(f"反馈增强：{feedback}")
            slides[idx]["speaker_notes"] = (slides[idx].get("speaker_notes", "") + f" | 反馈：{feedback}").strip()
        else:
            slides.append(
                {
                    "page": len(slides) + 1,
                    "title": "反馈优化补充页",
                    "bullets": [f"用户反馈：{feedback}", "本页根据反馈做定向增强。"],
                    "speaker_notes": "聚焦反馈点进行解释。",
                }
            )
        if scope.get("inject_kpi"):
            kpis = self._build_kpi_checklist(state.get("prompt", "")[:20] or "项目")
            slides.append(
                {
                    "page": len(slides) + 1,
                    "title": "可量化KPI指标清单（便于纳入区级考核）",
                    "bullets": [f"{k['name']}：目标 {k['target']}，周期 {k['cycle']}" for k in kpis],
                    "speaker_notes": "可直接纳入区级考核项。",
                    "kpis": kpis,
                }
            )
        prev_draft["slides"] = slides
        prev_draft["revision_note"] = feedback
        state["draft"] = prev_draft
        return state

    def _node_ppt_layout(self, state: Dict[str, Any]) -> Dict[str, Any]:
        draft = state.get("draft", {}) or {}
        slides = draft.get("slides", []) or []
        layout_slides: List[Dict[str, Any]] = []
        for idx, slide in enumerate(slides):
            bullets = [str(b)[:180] for b in (slide.get("bullets") or [])][:6]
            layout_slides.append(
                {
                    "page": idx + 1,
                    "title": str(slide.get("title", f"第 {idx+1} 页"))[:80],
                    "bullets": bullets,
                    "speaker_notes": str(slide.get("speaker_notes", ""))[:600],
                }
            )
        state["layout"] = {"slides": layout_slides}
        return state

    def _template_path(self, file_type: str, template_name: Optional[str] = None) -> Optional[Path]:
        if template_name:
            candidate = (self.template_dir / template_name).resolve()
            if candidate.exists() and str(candidate).startswith(str(self.template_dir)):
                return candidate
        default_name = "corporate_template.pptx" if file_type == "pptx" else "corporate_template.docx"
        candidate = (self.template_dir / default_name).resolve()
        return candidate if candidate.exists() else None

    @staticmethod
    def _docx_heading_level(heading: str) -> int:
        if re.search(r"[一二三四五六七八九十]+、", heading):
            return 1
        if re.search(r"\d+\.", heading):
            return 2
        return 1

    def _render_docx(self, draft: Dict[str, Any], target_path: Path, template_path: Optional[Path] = None):
        doc = Document(str(template_path)) if template_path else Document()
        title = draft.get("title") or "智能生成文档"
        doc.add_heading(str(title), level=0)

        analysis = draft.get("analysis", {}) or {}
        top_down = analysis.get("top_down", [])
        bottom_up = analysis.get("bottom_up", [])
        if top_down or bottom_up:
            doc.add_heading("上下文分析", level=1)
            if top_down:
                doc.add_paragraph("上位分解（Top-Down）：")
                for item in top_down:
                    doc.add_paragraph(str(item), style="List Bullet")
            if bottom_up:
                doc.add_paragraph("下位证据（Bottom-Up）：")
                for item in bottom_up:
                    doc.add_paragraph(str(item), style="List Bullet")

        for section in draft.get("sections", []) or []:
            heading = str(section.get("heading", "未命名章节"))
            doc.add_heading(heading, level=self._docx_heading_level(heading))
            for para in section.get("paragraphs", []) or []:
                doc.add_paragraph(str(para))

            kpis = section.get("kpis") or []
            if kpis:
                table = doc.add_table(rows=1, cols=5)
                hdr = table.rows[0].cells
                hdr[0].text = "指标"
                hdr[1].text = "口径"
                hdr[2].text = "目标"
                hdr[3].text = "周期"
                hdr[4].text = "责任人"
                for row in kpis:
                    cells = table.add_row().cells
                    cells[0].text = str(row.get("name", ""))
                    cells[1].text = str(row.get("formula", ""))
                    cells[2].text = str(row.get("target", ""))
                    cells[3].text = str(row.get("cycle", ""))
                    cells[4].text = str(row.get("owner", ""))

        target_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(target_path))

    def _render_pptx(self, layout: Dict[str, Any], target_path: Path, template_path: Optional[Path] = None):
        prs = Presentation(str(template_path)) if template_path else Presentation()
        slides = layout.get("slides", []) or []
        if not slides:
            slides = [{"title": "空白内容", "bullets": ["暂无内容"], "speaker_notes": ""}]

        for slide_data in slides:
            slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = str(slide_data.get("title", ""))
                if title_shape.text_frame and title_shape.text_frame.paragraphs:
                    p = title_shape.text_frame.paragraphs[0]
                    if p.runs:
                        p.runs[0].font.color.rgb = RGBColor(13, 89, 242)
                        p.runs[0].font.size = Pt(28)

            body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if body_shape and hasattr(body_shape, "text_frame"):
                tf = body_shape.text_frame
                tf.clear()
                bullets = slide_data.get("bullets", []) or []
                for idx, bullet in enumerate(bullets):
                    para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    para.text = str(bullet)
                    para.level = 0
                    if para.runs:
                        para.runs[0].font.size = Pt(18)

            # Corporate "logo" placeholder text
            logo_box = slide.shapes.add_textbox(Inches(11.2), Inches(0.2), Inches(1.0), Inches(0.35))
            logo_frame = logo_box.text_frame
            logo_frame.text = "猪你好运"
            if logo_frame.paragraphs and logo_frame.paragraphs[0].runs:
                run = logo_frame.paragraphs[0].runs[0]
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = RGBColor(13, 89, 242)

            notes = str(slide_data.get("speaker_notes", ""))
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        target_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(target_path))

    def _next_version(self, previous_job: Optional[Dict[str, Any]]) -> int:
        if not previous_job:
            return 1
        return int(previous_job.get("current_version", 1)) + 1

    def _node_renderer(self, state: Dict[str, Any]) -> Dict[str, Any]:
        previous_job = state.get("previous_job")
        req = state.get("requirements", {}) or {}
        file_type = req.get("file_type", "docx")
        prompt = state.get("prompt", "")
        session_id = state.get("session_id", "default")
        job_id = state.get("job_id") or (previous_job or {}).get("job_id") or uuid.uuid4().hex[:12]
        version = self._next_version(previous_job)

        safe_base = self._slug(prompt, fallback="workflow")
        filename = f"{safe_base}_{job_id}_v{version}.{file_type}"
        target_path = self.files_dir / filename
        template_path = self._template_path(file_type, req.get("template_name"))

        if file_type == "pptx":
            self._render_pptx(state.get("layout", {}) or state.get("draft", {}), target_path, template_path)
        else:
            self._render_docx(state.get("draft", {}), target_path, template_path)

        record = previous_job or {
            "job_id": job_id,
            "session_id": session_id,
            "prompt": prompt,
            "created_at": time.time(),
            "revision_history": [],
            "versions": [],
        }
        record["updated_at"] = time.time()
        record["current_version"] = version
        record["requirements"] = req
        record["plan"] = state.get("plan", {})
        record["draft"] = state.get("draft", {})
        record["layout"] = state.get("layout", {})
        record["file_type"] = file_type
        record["latest_filename"] = filename
        record["latest_download_url"] = f"/api/workflows/files/{filename}"
        record.setdefault("versions", []).append(
            {
                "version": version,
                "filename": filename,
                "download_url": f"/api/workflows/files/{filename}",
                "created_at": time.time(),
            }
        )
        if state.get("feedback"):
            record.setdefault("revision_history", []).append(
                {
                    "feedback": state.get("feedback"),
                    "scope": state.get("feedback_scope", {}),
                    "version": version,
                    "timestamp": time.time(),
                }
            )

        self._save_job(record)

        state["job_id"] = job_id
        state["version"] = version
        state["filename"] = filename
        state["download_url"] = f"/api/workflows/files/{filename}"
        state["result"] = record
        return state

    def _save_job(self, record: Dict[str, Any]):
        job_id = record.get("job_id")
        if not job_id:
            raise ValueError("job record missing job_id")
        path = self.jobs_dir / f"{job_id}.json"
        with path.open("w", encoding="utf-8") as f:
            json.dump(record, f, ensure_ascii=False, indent=2)

    def get_job(self, job_id: str) -> Optional[Dict[str, Any]]:
        path = self.jobs_dir / f"{job_id}.json"
        if not path.exists():
            return None
        with path.open("r", encoding="utf-8") as f:
            return json.load(f)

    def list_jobs(self, limit: int = 30) -> List[Dict[str, Any]]:
        records: List[Dict[str, Any]] = []
        for path in sorted(self.jobs_dir.glob("*.json"), key=lambda p: p.stat().st_mtime, reverse=True):
            try:
                with path.open("r", encoding="utf-8") as f:
                    record = json.load(f)
                records.append(
                    {
                        "job_id": record.get("job_id"),
                        "prompt": record.get("prompt"),
                        "file_type": record.get("file_type"),
                        "current_version": record.get("current_version"),
                        "latest_download_url": record.get("latest_download_url"),
                        "updated_at": record.get("updated_at"),
                    }
                )
                if len(records) >= limit:
                    break
            except Exception:
                continue
        return records

    def resolve_file_path(self, filename: str) -> Path:
        safe_name = Path(filename).name
        resolved = (self.files_dir / safe_name).resolve()
        if not str(resolved).startswith(str(self.files_dir)):
            raise ValueError("invalid filename")
        return resolved

    def generate(self, payload: Dict[str, Any]) -> Dict[str, Any]:
        state: Dict[str, Any] = {
            "mode": "generate",
            "job_id": uuid.uuid4().hex[:12],
            "session_id": payload.get("session_id", "default"),
            "prompt": payload.get("prompt", ""),
            "file_type": payload.get("file_type"),
            "target_words": payload.get("target_words"),
            "target_slides": payload.get("target_slides"),
            "template_name": payload.get("template_name"),
            "use_rag": payload.get("use_rag", True),
        }
        if HAS_LANGGRAPH and self.generate_graph:
            final_state = self.generate_graph.invoke(state)
        else:
            final_state = self._node_requirements_analyst(state)
            final_state = self._node_retriever(final_state)
            final_state = self._node_writer(final_state)
            if self._route_after_writer(final_state) == "ppt_layout":
                final_state = self._node_ppt_layout(final_state)
            final_state = self._node_renderer(final_state)

        return {
            "job_id": final_state.get("job_id"),
            "version": final_state.get("version"),
            "file_type": final_state.get("requirements", {}).get("file_type"),
            "download_url": final_state.get("download_url"),
            "filename": final_state.get("filename"),
            "requirements": final_state.get("requirements", {}),
            "analysis": (final_state.get("draft", {}) or {}).get("analysis", {}),
            "status": "completed",
        }

    def revise(self, payload: Dict[str, Any]) -> Dict[str, Any]:
        state: Dict[str, Any] = {
            "mode": "revise",
            "job_id": payload.get("job_id"),
            "session_id": payload.get("session_id", "default"),
            "feedback": payload.get("feedback", ""),
            "use_rag": True,
        }
        if HAS_LANGGRAPH and self.revise_graph:
            final_state = self.revise_graph.invoke(state)
        else:
            final_state = self._node_load_previous(state)
            final_state = self._node_human_feedback_router(final_state)
            final_state = self._node_retriever(final_state)
            final_state = self._node_rewrite(final_state)
            if self._route_after_writer(final_state) == "ppt_layout":
                final_state = self._node_ppt_layout(final_state)
            final_state = self._node_renderer(final_state)

        return {
            "job_id": final_state.get("job_id"),
            "version": final_state.get("version"),
            "file_type": final_state.get("requirements", {}).get("file_type"),
            "download_url": final_state.get("download_url"),
            "filename": final_state.get("filename"),
            "feedback_scope": final_state.get("feedback_scope", {}),
            "status": "completed",
        }
